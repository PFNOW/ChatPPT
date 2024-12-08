import os
from pptx import Presentation
from utils import remove_all_slides
from logger import LOG  # 引入日志模块
import pandas as pd  # 引入pandas模块读取表格数据
from pptx.util import Inches  # 引入 Inches 类，用于设置大小

# todo:表格功能支持除Excel以外的其他表格格式，如csv、txt、markdown、json等
# todo:支持用户自己选择的图表、插入组合图表
# todo:支持插入视频，但目前存在一些问题，如视频大小和位置的设置
# pptx库不支持对多媒体占位符的操作，直接插入视频的操作也有问题，必须指定视频的大小和位置，否则会报错

# 生成 PowerPoint 演示文稿
def generate_presentation(powerpoint_data, template_path: str, output_path: str):
    # 检查模板文件是否存在
    if not os.path.exists(template_path):
        LOG.error(f"模板文件 '{template_path}' 不存在。")  # 记录错误日志
        raise FileNotFoundError(f"模板文件 '{template_path}' 不存在。")

    prs = Presentation(template_path)  # 加载 PowerPoint 模板
    remove_all_slides(prs)  # 清除模板中的所有幻灯片
    prs.core_properties.title = powerpoint_data.title  # 设置 PowerPoint 的核心标题

    # 遍历所有幻灯片数据，生成对应的 PowerPoint 幻灯片
    for slide in powerpoint_data.slides:
        # 确保布局索引不超出范围，超出则使用默认布局
        if slide.layout_id >= len(prs.slide_layouts):
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[slide.layout_id]

        new_slide = prs.slides.add_slide(slide_layout)  # 添加新的幻灯片

        # 设置幻灯片标题
        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide.content.title
            LOG.debug(f"设置幻灯片标题: {slide.content.title}")

        # 添加文本内容
        for shape in new_slide.shapes:
            # 只处理非标题的文本框
            if shape.has_text_frame and not shape == new_slide.shapes.title:
                text_frame = shape.text_frame
                text_frame.clear()  # 清除原有内容
                # 将要点内容作为项目符号列表添加到文本框中
                for point in slide.content.bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0  # 项目符号的级别
                    LOG.debug(f"添加列表项: {point}")
                break

        # 插入图片
        if slide.content.image_path:
            image_full_path = os.path.join(os.getcwd(), slide.content.image_path)  # 构建图片的绝对路径
            if os.path.exists(image_full_path):
                # 插入图片到占位符中
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 18:  # 18 表示图片占位符
                        shape.insert_picture(image_full_path)
                        LOG.debug(f"插入图片: {image_full_path}")
                        break
            else:
                LOG.warning(f"图片路径 '{image_full_path}' 不存在，跳过此图片。")

        # 插入表格和图表
        if slide.content.table_path:
            table_full_path = os.path.join(os.getcwd(), slide.content.table_path)  # 构建表格的绝对路径
            if os.path.exists(table_full_path):
                # 插入表格到占位符中
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 12:  # 17 表示表格占位符
                        data = pd.read_excel(table_full_path) # 读取 Excel 文件
                        rows, cols = data.shape
                        table = shape.insert_table(rows+1, cols).table
                        # 设置表头
                        for col_num, column_title in enumerate(data.columns):
                            table.cell(0, col_num).text = str(column_title)
                        # 填充数据
                        for r in range(1, rows + 1):
                            for c in range(0, cols):
                                table.cell(r, c).text = str(data.iat[r - 1, c])
                        break
            else:
                LOG.warning(f"表格路径 '{table_full_path}' 不存在，跳过此表格。")

        # 插入图表
        if slide.content.chart_path:
            chart_full_path = os.path.join(os.getcwd(), slide.content.chart_path)  # 构建图表的绝对路径
            if os.path.exists(chart_full_path):
                # 插入图表到占位符中
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 12:  # 12 表示组合图表占位符
                        chart_type = slide.content.chart_type
                        chart_data = pd.read_excel(chart_full_path)
                        chart_data.plot(kind=chart_type, ax=shape.chart.chart_data.plot(kind=chart_type))
                        LOG.debug(f"插入图表: {chart_full_path}")
                        break
            else:
                LOG.warning(f"图表路径 '{chart_full_path}' 不存在，跳过此图表。")

        # 插入多媒体
        if slide.content.media_path:
            media_full_path = os.path.join(os.getcwd(), slide.content.media_path)  # 构建视频的绝对路径
            if os.path.exists(media_full_path):
                media_poster_path_full_path = ""  # 视频封面的路径
                # 构建视频封面的绝对路径
                if slide.content.media_poster_path:
                    media_poster_path_full_path = os.path.join(os.getcwd(), slide.content.media_poster_path)  # 构建视频封面的绝对路径
                if not os.path.exists(media_poster_path_full_path):
                    new_slide.shapes.add_movie(media_full_path, Inches(4.5), Inches(2.5), Inches(8), Inches(4.5))  # 插入媒体到幻灯片中
                else:
                    new_slide.shapes.add_movie(media_full_path, Inches(4.5), Inches(2.5), Inches(8), Inches(4.5), media_poster_path_full_path)  # 插入媒体到幻灯片中
                # 插入视频到占位符中
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 10:  # 10 表示视频占位符
                        LOG.debug(f"插入视频: {media_full_path}")
                        break
            else:
                LOG.warning(f"视频路径 '{media_full_path}' 不存在，跳过此视频。")

    # 保存生成的 PowerPoint 文件
    prs.save(output_path)
    LOG.info(f"演示文稿已保存到 '{output_path}'")
